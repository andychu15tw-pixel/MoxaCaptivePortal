from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ─── Color Palette ───────────────────────────────────────────────
C_NAVY   = RGBColor(0x1A, 0x3C, 0x6E)   # header bg
C_ORANGE = RGBColor(0xE8, 0x72, 0x2A)   # accent
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)   # content bg
C_DKGRAY = RGBColor(0x33, 0x33, 0x33)
C_MGRAY  = RGBColor(0x66, 0x66, 0x66)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_TEAL   = RGBColor(0x16, 0x7A, 0x8C)

W = Inches(13.33)   # 16:9 width
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ─── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=C_DKGRAY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def header_bar(slide, title, subtitle=None):
    """Full-width navy header bar at top."""
    add_rect(slide, 0, 0, W, Inches(1.35), fill=C_NAVY)
    add_text(slide, title, Inches(0.4), Inches(0.08), Inches(12), Inches(0.75),
             size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.78), Inches(12), Inches(0.45),
                 size=14, color=C_ORANGE)

def orange_accent(slide):
    """2px orange line under header."""
    add_rect(slide, 0, Inches(1.35), W, Inches(0.05), fill=C_ORANGE)

def bullet_box(slide, items, x, y, w, h, title=None, title_color=C_NAVY, bg=C_LGRAY):
    add_rect(slide, x, y, w, h, fill=bg)
    top = y
    if title:
        add_text(slide, title, x+Inches(0.15), top+Inches(0.1), w-Inches(0.3), Inches(0.4),
                 size=13, bold=True, color=title_color)
        top += Inches(0.42)
    for item in items:
        add_text(slide, "▸  " + item, x+Inches(0.15), top, w-Inches(0.3), Inches(0.38),
                 size=12, color=C_DKGRAY)
        top += Inches(0.36)

def code_box(slide, lines, x, y, w, h):
    add_rect(slide, x, y, w, h, fill=RGBColor(0x1E,0x1E,0x1E))
    top = y + Inches(0.12)
    for line in lines:
        add_text(slide, line, x+Inches(0.15), top, w-Inches(0.3), Inches(0.3),
                 size=10, color=RGBColor(0xAD,0xD8,0xE6), bold=False)
        top += Inches(0.27)

def tag(slide, text, x, y, w=None, bg=C_NAVY, fg=C_WHITE, size=11):
    if w is None: w = Inches(1.6)
    add_rect(slide, x, y, w, Inches(0.32), fill=bg)
    add_text(slide, text, x+Inches(0.05), y+Inches(0.03), w-Inches(0.1), Inches(0.28),
             size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)

def pass_tag(slide, x, y):
    tag(slide, "✔  PASS", x, y, Inches(1.3), bg=C_GREEN)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# full bg
add_rect(s, 0, 0, W, H, fill=C_NAVY)
# accent stripe
add_rect(s, 0, Inches(4.5), W, Inches(0.08), fill=C_ORANGE)

add_text(s, "Captive Portal Gateway", Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.1),
         size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s, "解決方案技術與操作說明", Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.8),
         size=30, color=C_ORANGE, align=PP_ALIGN.CENTER)
add_text(s, "Moxa V2400 / V3400  ·  Debian 12  ·  CoovaChilli + FreeRADIUS + daloRADIUS",
         Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.55),
         size=16, color=RGBColor(0xB0,0xC8,0xE8), align=PP_ALIGN.CENTER)

add_text(s, "FAE / AE 技術參考文件", Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4),
         size=13, color=RGBColor(0x80,0xA0,0xC0), align=PP_ALIGN.CENTER, italic=True)
add_text(s, "2026-05", Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.35),
         size=12, color=RGBColor(0x60,0x80,0xA0), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 議程
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "議程 Agenda")
orange_accent(s)

chapters = [
    ("01", "解決方案定位與硬體平台"),
    ("02", "軟體架構總覽"),
    ("03", "系統網路拓樸"),
    ("04", "CoovaChilli 核心機制（tun0 / L2）"),
    ("05", "Client 認證流程"),
    ("06", "WISPr 限速機制與調整"),
    ("07", "daloRADIUS 管理後台操作"),
    ("08", "CoA 即時踢人"),
    ("09", "v1 MVP 驗證結果"),
    ("10", "後續 v2 規劃"),
]

col_x = [Inches(0.5), Inches(7.0)]
for i, (num, title) in enumerate(chapters):
    col = i // 5
    row = i % 5
    x = col_x[col]
    y = Inches(1.6) + row * Inches(1.05)
    add_rect(s, x, y, Inches(5.8), Inches(0.85), fill=C_LGRAY)
    add_rect(s, x, y, Inches(0.6), Inches(0.85), fill=C_NAVY)
    add_text(s, num, x+Inches(0.02), y+Inches(0.18), Inches(0.56), Inches(0.5),
             size=13, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, title, x+Inches(0.72), y+Inches(0.18), Inches(4.9), Inches(0.5),
             size=13, color=C_DKGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 解決方案定位
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "解決方案定位", "Network Access Control + Captive Portal + AAA on Moxa Industrial PC")
orange_accent(s)

cards = [
    (C_NAVY,   "🎯  使用情境",
     ["工廠訪客網路管控", "飯店/辦公室 Wi-Fi 認證", "工業現場 NAC 邊緣閘道", "IoT 設備網路存取控制"]),
    (C_TEAL,   "⚙️  核心功能",
     ["Captive Portal 強制登入", "RADIUS AAA 認證/授權/計費", "Per-user 頻寬限制", "即時踢人（CoA / PoD）"]),
    (C_ORANGE, "💡  設計原則",
     ["能用現成套件就用", "架構越單純越好", "自寫 code ≤ 100 行 shell", "可示範、可壓測、可上線"]),
]

for i, (color, title, items) in enumerate(cards):
    x = Inches(0.35) + i * Inches(4.3)
    add_rect(s, x, Inches(1.55), Inches(4.1), Inches(5.65), fill=C_LGRAY)
    add_rect(s, x, Inches(1.55), Inches(4.1), Inches(0.52), fill=color)
    add_text(s, title, x+Inches(0.12), Inches(1.6), Inches(3.9), Inches(0.44),
             size=14, bold=True, color=C_WHITE)
    top = Inches(2.25)
    for item in items:
        add_text(s, "▸  " + item, x+Inches(0.18), top, Inches(3.8), Inches(0.42),
                 size=12.5, color=C_DKGRAY)
        top += Inches(0.55)

# HW badge
add_rect(s, Inches(0.35), Inches(6.9), Inches(12.6), Inches(0.38), fill=C_NAVY)
add_text(s, "硬體平台：Moxa V2400 / V3400  x86 工業電腦  |  OS：Debian 12 (bookworm)  |  架構：aarch64 / x86",
         Inches(0.5), Inches(6.93), Inches(12.3), Inches(0.34),
         size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 軟體架構總覽
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "軟體架構總覽", "全套現成套件 — 零自寫應用程式碼")
orange_accent(s)

rows = [
    ("Captive Portal", "CoovaChilli 1.6", "coova-chilli", "UAM redirect、DHCP、DNS hijack、RADIUS client、leaky bucket 限速"),
    ("AAA", "FreeRADIUS 3", "freeradius + freeradius-mysql", "認證 / 授權 / Accounting，後端 MariaDB"),
    ("Admin Web UI", "daloRADIUS 1.3", "GitHub release (PHP)", "使用者管理 / 報表 / CoA 踢人"),
    ("Web Server", "Apache 2.4 + PHP 8.2", "apache2, php, php-mysql…", "daloRADIUS 執行環境 + CGI portal 頁"),
    ("Database", "MariaDB", "mariadb-server", "FreeRADIUS + daloRADIUS 共用同一 instance"),
    ("Firewall/NAT", "nftables", "nftables", "Pre/post-auth 規則、masquerade"),
    ("Cellular", "ModemManager", "modemmanager", "4G/5G 模組管理（v2）"),
    ("WAN Failover", "keepalived", "keepalived", "ping healthcheck + route 切換（v2）"),
]

col_w = [Inches(1.7), Inches(2.0), Inches(2.5), Inches(6.5)]
col_x2 = [Inches(0.25), Inches(1.97), Inches(4.0), Inches(6.55)]
headers = ["層", "套件", "Debian Package", "用途"]

# header row
add_rect(s, Inches(0.25), Inches(1.55), Inches(12.8), Inches(0.42), fill=C_NAVY)
for j, h in enumerate(headers):
    add_text(s, h, col_x2[j]+Inches(0.05), Inches(1.57), col_w[j], Inches(0.38),
             size=12, bold=True, color=C_WHITE)

for i, row in enumerate(rows):
    y = Inches(1.97) + i * Inches(0.6)
    bg = C_LGRAY if i % 2 == 0 else C_WHITE
    add_rect(s, Inches(0.25), y, Inches(12.8), Inches(0.58), fill=bg)
    colors = [C_NAVY, C_TEAL, C_DKGRAY, C_MGRAY]
    for j, cell in enumerate(row):
        sz = 11 if j == 3 else 12
        bold = j in (0, 1)
        c = colors[j]
        add_text(s, cell, col_x2[j]+Inches(0.05), y+Inches(0.1), col_w[j]-Inches(0.1), Inches(0.42),
                 size=sz, bold=bold, color=c)

add_text(s, "★  自寫 code 量目標：≤ 100 行 shell（cellular failover healthcheck），其餘全靠設定檔",
         Inches(0.3), Inches(6.95), Inches(12.5), Inches(0.38),
         size=12, bold=True, color=C_ORANGE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 系統網路拓樸
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "系統網路拓樸", "物理與邏輯介面配置")
orange_accent(s)

# Draw topology boxes
def topo_box(slide, label, sub, x, y, w=Inches(2.2), h=Inches(0.9), bg=C_NAVY, fg=C_WHITE):
    add_rect(slide, x, y, w, h, fill=bg)
    add_text(slide, label, x+Inches(0.1), y+Inches(0.05), w-Inches(0.2), Inches(0.42),
             size=13, bold=True, color=fg, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, sub, x+Inches(0.1), y+Inches(0.45), w-Inches(0.2), Inches(0.35),
                 size=10, color=RGBColor(0xB0,0xD0,0xF0), align=PP_ALIGN.CENTER)

def arrow_h(slide, x1, y, x2, color=C_NAVY, label=None):
    line = slide.shapes.add_connector(1, x1, y, x2, y)
    line.line.color.rgb = color
    line.line.width = Pt(2.0)
    if label:
        mx = (x1+x2)//2
        add_text(slide, label, mx-Inches(0.6), y-Inches(0.32), Inches(1.2), Inches(0.28),
                 size=9, color=color, align=PP_ALIGN.CENTER, italic=True)

def arrow_v(slide, x, y1, y2, color=C_NAVY, label=None):
    line = slide.shapes.add_connector(1, x, y1, x, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.0)
    if label:
        my = (y1+y2)//2
        add_text(slide, label, x+Inches(0.05), my-Inches(0.15), Inches(1.5), Inches(0.28),
                 size=9, color=color, italic=True)

# Internet
topo_box(s, "🌐  Internet", None, Inches(11.0), Inches(2.3), bg=C_TEAL)

# Moxa box
add_rect(s, Inches(3.8), Inches(1.5), Inches(5.6), Inches(5.5),
         fill=RGBColor(0xEE,0xF5,0xFF), line=C_NAVY, line_w=Pt(2))
add_text(s, "Moxa V2400 / V3400  —  Debian 12", Inches(3.9), Inches(1.55), Inches(5.4), Inches(0.38),
         size=11, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

# Internal services
services = [
    ("CoovaChilli", "tun0 192.168.182.0/24"),
    ("FreeRADIUS",  "127.0.0.1:1812/1813"),
    ("MariaDB",     "radius DB"),
    ("Apache+PHP",  "daloRADIUS /daloradius"),
]
for i, (svc, note) in enumerate(services):
    sx = Inches(4.0) + (i % 2) * Inches(2.7)
    sy = Inches(2.1) + (i // 2) * Inches(1.1)
    add_rect(s, sx, sy, Inches(2.5), Inches(0.82), fill=C_NAVY)
    add_text(s, svc,  sx+Inches(0.1), sy+Inches(0.05), Inches(2.3), Inches(0.38),
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, note, sx+Inches(0.1), sy+Inches(0.42), Inches(2.3), Inches(0.32),
             size=9, color=C_ORANGE, align=PP_ALIGN.CENTER)

# eth0 label
add_rect(s, Inches(9.4), Inches(2.9), Inches(1.6), Inches(0.38), fill=C_ORANGE)
add_text(s, "eth0 (WAN)", Inches(9.42), Inches(2.92), Inches(1.56), Inches(0.35),
         size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
arrow_h(s, Inches(11.0), Inches(3.1), Inches(9.4), color=C_ORANGE, label="Internet")

# eth1 label
add_rect(s, Inches(2.2), Inches(2.9), Inches(1.6), Inches(0.38), fill=C_TEAL)
add_text(s, "eth1 (LAN)", Inches(2.22), Inches(2.92), Inches(1.56), Inches(0.35),
         size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Wi-Fi AP
topo_box(s, "📶  Wi-Fi AP", "Bridge Mode (L2)", Inches(0.2), Inches(2.2), bg=C_TEAL)
arrow_h(s, Inches(2.4), Inches(2.65), Inches(3.8), color=C_TEAL, label="L2 Bridge")

# wwan0
add_rect(s, Inches(9.4), Inches(4.2), Inches(1.6), Inches(0.38), fill=RGBColor(0x8E,0x44,0xAD))
add_text(s, "wwan0 (Cellular)", Inches(9.42), Inches(4.22), Inches(1.56), Inches(0.35),
         size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s, "（v2）", Inches(9.5), Inches(4.6), Inches(1.4), Inches(0.28),
         size=9, color=RGBColor(0x8E,0x44,0xAD), align=PP_ALIGN.CENTER, italic=True)

# Mgmt
topo_box(s, "💻  Admin PC", "10.90.x.x / MGMT", Inches(0.2), Inches(5.0),
         bg=RGBColor(0x27,0x6B,0x45), w=Inches(2.2))
add_text(s, "daloRADIUS UI\nhttps://10.90.35.42/daloradius",
         Inches(0.2), Inches(5.95), Inches(2.2), Inches(0.45),
         size=9, color=RGBColor(0x27,0x6B,0x45))

# Client
topo_box(s, "📱  Client", "192.168.182.x", Inches(0.2), Inches(3.8),
         bg=RGBColor(0x9B,0x59,0xB6))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CoovaChilli 核心機制
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "CoovaChilli 核心機制", "為什麼需要 tun0？Wi-Fi AP 為何要 Bridge Mode？")
orange_accent(s)

# Left: tun0 explanation
add_rect(s, Inches(0.3), Inches(1.55), Inches(6.0), Inches(5.65), fill=C_LGRAY)
add_rect(s, Inches(0.3), Inches(1.55), Inches(6.0), Inches(0.42), fill=C_NAVY)
add_text(s, "① 為什麼需要 tun0？", Inches(0.4), Inches(1.58), Inches(5.8), Inches(0.38),
         size=13, bold=True, color=C_WHITE)

flow_items = [
    (C_TEAL,   "Client 封包",      "eth1 raw socket → CoovaChilli userspace"),
    (C_ORANGE, "chilli 判斷",      "未認證 → redirect portal / 已認證 → 放行"),
    (C_NAVY,   "inject tun0",      "chilli 把封包寫入 tun0（虛擬介面）"),
    (C_GREEN,  "kernel 接手",      "routing / nftables / masquerade"),
    (C_TEAL,   "eth0 → Internet",  "封包送出 WAN"),
]
top = Inches(2.2)
for bg, label, desc in flow_items:
    add_rect(s, Inches(0.5), top, Inches(1.5), Inches(0.52), fill=bg)
    add_text(s, label, Inches(0.52), top+Inches(0.1), Inches(1.46), Inches(0.38),
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, Inches(2.15), top+Inches(0.1), Inches(3.9), Inches(0.38),
             size=11, color=C_DKGRAY)
    if top < Inches(4.3):
        add_text(s, "↓", Inches(1.15), top+Inches(0.52), Inches(0.4), Inches(0.28),
                 size=14, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    top += Inches(0.82)

add_text(s, "eth0/eth1 寫入 = 送上實體網線\ntun0 寫入   = 注入 kernel IP 堆疊 → routing 接手",
         Inches(0.45), Inches(6.5), Inches(5.8), Inches(0.62),
         size=11, color=C_NAVY, bold=True)

# Right: L2 Bridge Mode
add_rect(s, Inches(6.7), Inches(1.55), Inches(6.3), Inches(5.65), fill=C_LGRAY)
add_rect(s, Inches(6.7), Inches(1.55), Inches(6.3), Inches(0.42), fill=C_TEAL)
add_text(s, "② Wi-Fi AP 必須 Bridge Mode（L2）", Inches(6.8), Inches(1.58), Inches(6.1), Inches(0.38),
         size=13, bold=True, color=C_WHITE)

# Table comparison
modes = [
    ("AP 模式", "chilli 看到的 MAC", "結果", True),
    ("Bridge (L2) ✔", "每個 client 真實 MAC", "✅ 各自認證/計費", False),
    ("Router/NAT ✗", "只有 AP 的一個 MAC", "❌ 所有人同一 session", False),
]
ty = Inches(2.15)
for i, (m, mac, result, is_header) in enumerate(modes):
    bg = C_NAVY if is_header else (RGBColor(0xD5,0xF5,0xE3) if "✔" in m else RGBColor(0xFD,0xED,0xEC))
    add_rect(s, Inches(6.8), ty, Inches(6.0), Inches(0.5), fill=bg)
    fg = C_WHITE if is_header else C_DKGRAY
    add_text(s, m,      Inches(6.85), ty+Inches(0.08), Inches(1.6),  Inches(0.38), size=11, bold=is_header, color=fg)
    add_text(s, mac,    Inches(8.5),  ty+Inches(0.08), Inches(2.3),  Inches(0.38), size=10, bold=is_header, color=fg)
    add_text(s, result, Inches(10.85),ty+Inches(0.08), Inches(1.85), Inches(0.38), size=10, bold=is_header, color=fg)
    ty += Inches(0.52)

add_text(s, "chilli 用 MAC 識別每個 client：\n▸  哪個 MAC → 是否通過 portal 認證\n▸  哪個 MAC → 對應哪個 RADIUS session\n▸  L2 bridge 讓 client MAC 直接穿透 AP 到 eth1",
         Inches(6.85), Inches(3.4), Inches(5.9), Inches(1.7),
         size=12, color=C_DKGRAY)

add_text(s, "AP 設定重點：關掉 DHCP、關掉 NAT、選 Access Point / Bridge 模式\nDHCP 由 chilli 統一發放（192.168.182.0/24）",
         Inches(6.85), Inches(5.3), Inches(5.9), Inches(0.9),
         size=11, bold=True, color=C_NAVY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Client 認證流程
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "Client 認證流程", "從接入到上網的完整步驟")
orange_accent(s)

steps = [
    (C_TEAL,   "① DHCP",       "Client 連 eth1 → chilli 發 192.168.182.x IP\nDNS 指向 192.168.182.1（chilli 劫持）"),
    (C_ORANGE, "② HTTP 攔截",  "Client 開任意 http:// → chilli 回 302\n導向 http://192.168.182.1/cgi-bin/hotspotlogin.cgi"),
    (C_NAVY,   "③ Portal 登入","User 輸入帳密 → CHAP 加密 → hotspotlogin.cgi\n送 UAM redirect 給 chilli"),
    (C_TEAL,   "④ RADIUS Auth","chilli → FreeRADIUS Access-Request\nFreeRADIUS 查 MariaDB radcheck 表"),
    (C_GREEN,  "⑤ Access-Accept","FreeRADIUS 回傳 Access-Accept + reply attrs\n（Session-Timeout、WISPr-Bandwidth、…）"),
    (C_ORANGE, "⑥ 上網",       "chilli 設定 leaky bucket 限速 + 開放轉發\nClient 流量 tun0 → kernel → NAT → Internet"),
]

col_w2 = Inches(6.2)
for i, (bg, title, desc) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = Inches(0.3) + col * Inches(6.5)
    y = Inches(1.6) + row * Inches(1.75)
    add_rect(s, x, y, col_w2, Inches(1.6), fill=C_LGRAY)
    add_rect(s, x, y, Inches(0.6), Inches(1.6), fill=bg)
    add_text(s, str(i+1), x+Inches(0.05), y+Inches(0.55), Inches(0.5), Inches(0.5),
             size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x+Inches(0.72), y+Inches(0.12), Inches(5.3), Inches(0.42),
             size=13, bold=True, color=bg)
    add_text(s, desc,  x+Inches(0.72), y+Inches(0.55), Inches(5.3), Inches(0.92),
             size=11, color=C_DKGRAY)

add_text(s, "OS Captive Detection：iOS / Android / Windows 開機自動偵測 HTTP 攔截 → 彈出 CNA 視窗（零點擊）",
         Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.36),
         size=11, bold=True, color=C_NAVY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7B — 認證流程序列圖 (Sequence Diagram)
# ════════════════════════════════════════════════════════════════════════════
from pptx.oxml.ns import qn
from lxml import etree

def add_arrow(slide, x1, y1, x2, y2, color=C_NAVY, width_pt=1.6, dashed=False, head_back=False):
    """Add line with arrowhead at end (or both if head_back=True)."""
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    spPr = line._element.find(qn('p:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    if dashed:
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle'); tailEnd.set('w', 'med'); tailEnd.set('len', 'med')
    if head_back:
        headEnd = etree.SubElement(ln, qn('a:headEnd'))
        headEnd.set('type', 'triangle'); headEnd.set('w', 'med'); headEnd.set('len', 'med')
    return line

s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "認證流程序列圖", "對應上頁 6 步驟的時序視圖")
orange_accent(s)

# ── Lifeline columns (4 lanes, matching Slide 7 mention) ──
lanes = [
    ("📱  Client",     Inches(1.6),  C_TEAL),
    ("🖧  chilli",     Inches(5.2),  C_NAVY),
    ("🛡  FreeRADIUS", Inches(8.8),  C_ORANGE),
    ("🗄  MariaDB",    Inches(11.9), RGBColor(0x6C,0x3A,0x8C)),
]
lane_top = Inches(1.5)
lane_bottom = Inches(6.8)
for label, cx, color in lanes:
    add_rect(s, cx - Inches(1.0), lane_top, Inches(2.0), Inches(0.5), fill=color)
    add_text(s, label, cx - Inches(0.98), lane_top + Inches(0.06),
             Inches(1.96), Inches(0.4),
             size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # dashed lifeline
    lf = s.shapes.add_connector(1, cx, lane_top + Inches(0.55), cx, lane_bottom)
    lf.line.color.rgb = C_MGRAY
    lf.line.width = Pt(0.75)
    lnel = lf._element.find(qn('p:spPr')).find(qn('a:ln'))
    if lnel is None:
        lnel = etree.SubElement(lf._element.find(qn('p:spPr')), qn('a:ln'))
    prst = etree.SubElement(lnel, qn('a:prstDash'))
    prst.set('val', 'dash')

LX_CLIENT = Inches(1.6)
LX_CHILLI = Inches(5.2)
LX_RADIUS = Inches(8.8)
LX_DB     = Inches(11.9)

# ── 6 大步驟（與 Slide 7 對齊）每步驟可含 1-2 條箭頭 ──
# group: (num, color, group_label, [(x1, x2, sub_label, dashed), ...])
groups = [
    ("①", C_TEAL,   "DHCP",
        [(LX_CLIENT, LX_CHILLI, "Discover / Offer / Request / ACK  →  192.168.182.x + DNS=192.168.182.1", False)]),
    ("②", C_ORANGE, "HTTP 攔截",
        [(LX_CLIENT, LX_CHILLI, "HTTP GET  http://任意網站", False),
         (LX_CHILLI, LX_CLIENT, "302 Redirect  →  /cgi-bin/hotspotlogin.cgi", True)]),
    ("③", C_NAVY,   "Portal 登入",
        [(LX_CLIENT, LX_CHILLI, "POST 帳號/密碼  →  CGI 計算 CHAP-Response", False)]),
    ("④", C_TEAL,   "RADIUS Auth",
        [(LX_CHILLI, LX_RADIUS, "Access-Request (User-Name + CHAP-Password)", False),
         (LX_RADIUS, LX_DB,     "SELECT radcheck WHERE username=", False)]),
    ("⑤", C_GREEN,  "Access-Accept",
        [(LX_DB,     LX_RADIUS, "rows (Cleartext-Password + WISPr attrs)", True),
         (LX_RADIUS, LX_CHILLI, "Access-Accept  +  Session-Timeout / WISPr-Bandwidth", True)]),
    ("⑥", C_GREEN,  "上網",
        [(LX_CHILLI, LX_CLIENT, "授權 ACL + leaky bucket 限速 + 302 google.com", True)]),
]

# Render
y_cursor = 2.0
for num, color, gtitle, arrows in groups:
    # Group label badge on far left
    badge_y = Inches(y_cursor)
    add_rect(s, Inches(0.15), badge_y, Inches(1.05), Inches(0.4), fill=color)
    add_text(s, f"{num} {gtitle}", Inches(0.15), badge_y + Inches(0.06),
             Inches(1.05), Inches(0.32),
             size=10.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Draw arrows for this group
    for j, (x1, x2, label, dashed) in enumerate(arrows):
        ay = Inches(y_cursor + 0.12 + j * 0.42)
        add_arrow(s, x1, ay, x2, ay, color=color, dashed=dashed, width_pt=1.8)
        mx = (x1 + x2) // 2
        add_text(s, label, mx - Inches(2.4), ay - Inches(0.32),
                 Inches(4.8), Inches(0.28),
                 size=9.5, color=color, align=PP_ALIGN.CENTER)

    # Advance cursor by number of arrows
    y_cursor += 0.45 * max(1, len(arrows)) + 0.25

# ── Footer ──
add_rect(s, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.42),
         fill=C_LGRAY)
add_text(s,
         "─── = Request  ┄┄ = Response   |   ④⑤ 走 RADIUS UDP 1812/1813   |   ⑥ 後：interim acct 每 300s",
         Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35),
         size=10, color=C_DKGRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WISPr 限速
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "WISPr 頻寬限制", "Per-user 限速機制、調整方式與驗證")
orange_accent(s)

# Left: mechanism
add_rect(s, Inches(0.3), Inches(1.55), Inches(5.5), Inches(5.65), fill=C_LGRAY)
add_rect(s, Inches(0.3), Inches(1.55), Inches(5.5), Inches(0.42), fill=C_NAVY)
add_text(s, "機制說明", Inches(0.4), Inches(1.58), Inches(5.3), Inches(0.38),
         size=13, bold=True, color=C_WHITE)

mech = [
    "RADIUS reply 回傳 WISPr 屬性",
    "WISPr-Bandwidth-Max-Down（bps）",
    "WISPr-Bandwidth-Max-Up（bps）",
    "",
    "CoovaChilli 1.6 用 Leaky Bucket 實作",
    "（非 tc/HTB，userspace 軟體限速）",
    "",
    "適用範圍：tun0 所有轉發流量",
    "（client ↔ internet）",
]
top = Inches(2.1)
for m in mech:
    if m:
        prefix = "▸  " if not m.startswith("（") else "      "
        add_text(s, prefix + m, Inches(0.45), top, Inches(5.2), Inches(0.36),
                 size=11, color=C_DKGRAY if not m.startswith("WISPr") else C_NAVY,
                 bold=m.startswith("WISPr"))
    top += Inches(0.36)

# Right: table + test result
add_rect(s, Inches(6.1), Inches(1.55), Inches(6.9), Inches(2.5), fill=C_LGRAY)
add_rect(s, Inches(6.1), Inches(1.55), Inches(6.9), Inches(0.42), fill=C_TEAL)
add_text(s, "限速值對照表", Inches(6.2), Inches(1.58), Inches(6.7), Inches(0.38),
         size=13, bold=True, color=C_WHITE)

speed_rows = [
    ("512 Kbps", "512000", "~64 KB/s", "1MB 約 15.6s"),
    ("1 Mbps",   "1000000", "~125 KB/s", "1MB 約 8.0s ← 實測"),
    ("5 Mbps",   "5000000", "~625 KB/s", "1MB 約 1.6s"),
    ("無限制",   "0 或刪row", "—", "—"),
]
ty2 = Inches(2.05)
cols_x = [Inches(6.15), Inches(7.8), Inches(9.4), Inches(11.0)]
cols_w2 = [Inches(1.6), Inches(1.5), Inches(1.6), Inches(2.0)]
hd = ["速度", "填入值 (bps)", "實際速率", "測試結果"]
add_rect(s, Inches(6.1), ty2, Inches(6.9), Inches(0.4), fill=C_NAVY)
for j, h in enumerate(hd):
    add_text(s, h, cols_x[j]+Inches(0.05), ty2+Inches(0.05), cols_w2[j], Inches(0.32),
             size=10, bold=True, color=C_WHITE)
ty2 += Inches(0.42)
for i, row in enumerate(speed_rows):
    bg = RGBColor(0xD5,0xF5,0xE3) if "實測" in row[3] else (C_LGRAY if i%2==0 else C_WHITE)
    add_rect(s, Inches(6.1), ty2, Inches(6.9), Inches(0.4), fill=bg)
    for j, cell in enumerate(row):
        add_text(s, cell, cols_x[j]+Inches(0.05), ty2+Inches(0.05), cols_w2[j], Inches(0.32),
                 size=10, color=C_DKGRAY if "實測" not in row[3] else C_NAVY,
                 bold="實測" in row[3])
    ty2 += Inches(0.42)

# Adjust procedure
add_rect(s, Inches(6.1), Inches(4.2), Inches(6.9), Inches(2.98), fill=C_LGRAY)
add_rect(s, Inches(6.1), Inches(4.2), Inches(6.9), Inches(0.42), fill=C_ORANGE)
add_text(s, "調整與測試流程", Inches(6.2), Inches(4.23), Inches(6.7), Inches(0.38),
         size=13, bold=True, color=C_WHITE)

adjust_steps = [
    "1. daloRADIUS → testuser → Reply Attributes → 改值",
    "   或 SQL: UPDATE radreply SET value=5000000 WHERE ...",
    "2. Disconnect User（CoA 踢人）",
    "3. Client 重新登入 portal（套新值）",
    "4. wget speedtest.tele2.net/1MB.zip -O /dev/null",
    "5. 驗算：預期秒數 = 1048576 ÷ (bps÷8)",
]
top2 = Inches(4.72)
for step in adjust_steps:
    bold = step.startswith(("1.","2.","3.","4.","5."))
    add_text(s, step, Inches(6.25), top2, Inches(6.6), Inches(0.36),
             size=10.5, color=C_DKGRAY, bold=bold)
    top2 += Inches(0.38)

add_text(s, "⚠ 修改 radreply 後必須斷線重連，chilli 在新 auth 時才讀新值",
         Inches(6.15), Inches(6.9), Inches(6.8), Inches(0.38),
         size=11, bold=True, color=C_ORANGE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — daloRADIUS 操作
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "daloRADIUS 管理後台", "URL / 帳號 / 主要操作功能")
orange_accent(s)

# URL info
add_rect(s, Inches(0.3), Inches(1.6), Inches(12.7), Inches(0.55), fill=C_NAVY)
add_text(s, "🔗  https://<GW-IP>/daloradius/     帳號: administrator / radius（預設，建議上線前改）",
         Inches(0.5), Inches(1.67), Inches(12.3), Inches(0.42),
         size=13, bold=True, color=C_ORANGE)

ops = [
    ("👤  使用者管理", "Management → Users",
     ["新增帳號：New User / New User Quick Add",
      "編輯帳號：Edit User → Account Info / Reply Attributes",
      "搜尋帳號：Search Users",
      "刪除帳號：Remove Users"]),
    ("📊  即時監控", "Reports → Online Users / Last Connections",
     ["查看目前在線 client 清單",
      "查看最近認證記錄",
      "Session 開始時間 / IP / MAC",
      "Accounting 流量統計"]),
    ("🔧  CoA 踢人", "Config → Maintenance → Disconnect User",
     ["選 username → Packet Type: PoD",
      "NAS IP: chilli 127.0.0.1",
      "NAS Port: 3799",
      "點 Disconnect User → 收到 Disconnect-ACK"]),
    ("📋  Accounting 報表", "Accounting → All / By Username",
     ["查詢每個 user 的用量",
      "起訖時間 / 上下行流量",
      "acctterminatecause 顯示斷線原因",
      "Admin-Reset = 被 CoA 踢掉"]),
]

for i, (title, path, items) in enumerate(ops):
    col = i % 2
    row = i // 2
    x = Inches(0.3) + col * Inches(6.5)
    y = Inches(2.35) + row * Inches(2.55)
    add_rect(s, x, y, Inches(6.2), Inches(2.4), fill=C_LGRAY)
    add_rect(s, x, y, Inches(6.2), Inches(0.5), fill=C_TEAL)
    add_text(s, title, x+Inches(0.12), y+Inches(0.05), Inches(5.0), Inches(0.42),
             size=13, bold=True, color=C_WHITE)
    add_text(s, path,  x+Inches(0.12), y+Inches(0.52), Inches(6.0), Inches(0.32),
             size=11, italic=True, color=C_NAVY)
    tp = y + Inches(0.9)
    for item in items:
        add_text(s, "▸  " + item, x+Inches(0.2), tp, Inches(5.8), Inches(0.36),
                 size=11, color=C_DKGRAY)
        tp += Inches(0.36)

add_text(s, "daloRADIUS 目錄結構（v1.3）：所有 PHP 直接在 /opt/daloradius/，無 app/operators/ 子目錄",
         Inches(0.3), Inches(7.08), Inches(12.7), Inches(0.32),
         size=10, color=C_MGRAY, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CoA 踢人
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "CoA / PoD 即時踢人", "Change of Authorization — 強制終止 client session")
orange_accent(s)

# Left: flow
add_rect(s, Inches(0.3), Inches(1.55), Inches(5.8), Inches(5.65), fill=C_LGRAY)
add_rect(s, Inches(0.3), Inches(1.55), Inches(5.8), Inches(0.42), fill=C_NAVY)
add_text(s, "PoD 封包流向", Inches(0.4), Inches(1.58), Inches(5.6), Inches(0.38),
         size=13, bold=True, color=C_WHITE)

pod_steps = [
    (C_TEAL,   "Admin 按下\nDisconnect User", "daloRADIUS Web UI"),
    (C_ORANGE, "radclient 送出\nDisconnect-Request", "UDP 127.0.0.1:3799\nUser-Name = testuser"),
    (C_NAVY,   "CoovaChilli\n收到 PoD", "coaport=3799\n驗證 radiussecret"),
    (C_GREEN,  "Disconnect-ACK\n回傳", "chilli 踢掉 client session\nleaky bucket 釋放"),
    (C_TEAL,   "radacct 更新", "AcctStopTime = 現在\nAcctTerminateCause = Admin-Reset"),
]
top3 = Inches(2.1)
for j, (bg, label, note) in enumerate(pod_steps):
    add_rect(s, Inches(0.5), top3, Inches(2.2), Inches(0.72), fill=bg)
    add_text(s, label, Inches(0.52), top3+Inches(0.05), Inches(2.16), Inches(0.65),
             size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, note, Inches(2.85), top3+Inches(0.1), Inches(3.1), Inches(0.55),
             size=10, color=C_DKGRAY)
    if j < len(pod_steps)-1:
        add_text(s, "↓", Inches(1.5), top3+Inches(0.72), Inches(0.4), Inches(0.22),
                 size=12, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    top3 += Inches(0.95)

# Right: verification
add_rect(s, Inches(6.4), Inches(1.55), Inches(6.6), Inches(2.5), fill=C_LGRAY)
add_rect(s, Inches(6.4), Inches(1.55), Inches(6.6), Inches(0.42), fill=C_TEAL)
add_text(s, "實際執行結果", Inches(6.5), Inches(1.58), Inches(6.4), Inches(0.38),
         size=13, bold=True, color=C_WHITE)

code_box(s, [
    "Sent Disconnect-Request Id 237",
    "  from 0.0.0.0:39556",
    "  to 127.0.0.1:3799 length 30",
    "User-Name = \"testuser\"",
    "Received Disconnect-ACK Id 237",
    "  from 127.0.0.1:3799",
    "  to 127.0.0.1:39556 length 20",
], Inches(6.4), Inches(2.0), Inches(6.6), Inches(2.0))

# DB verification
add_rect(s, Inches(6.4), Inches(4.15), Inches(6.6), Inches(3.05), fill=C_LGRAY)
add_rect(s, Inches(6.4), Inches(4.15), Inches(6.6), Inches(0.42), fill=C_NAVY)
add_text(s, "radacct 驗證", Inches(6.5), Inches(4.18), Inches(6.4), Inches(0.38),
         size=13, bold=True, color=C_WHITE)

code_box(s, [
    "mysql> SELECT username,",
    "  acctstarttime, acctstoptime,",
    "  acctterminatecause",
    "  FROM radacct",
    "  WHERE username='testuser'",
    "  ORDER BY acctstarttime DESC",
    "  LIMIT 1;",
    "",
    "testuser | 08:45:09 | 09:22:09",
    "acctterminatecause: Admin-Reset ✔",
], Inches(6.4), Inches(4.6), Inches(6.6), Inches(2.55))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — v1 MVP 驗證結果
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "v1 MVP 驗證結果", "11 項測試項目全數通過")
orange_accent(s)

tests = [
    ("DHCP",            "Client 接 LAN → 取得 192.168.182.x IP"),
    ("Portal Redirect", "HTTP 任意站 → 302 導向 hotspotlogin.cgi"),
    ("OS 偵測 (CNA)",   "iOS / Android / Windows 自動彈認證視窗"),
    ("RADIUS 認證",     "testuser/test1234 登入成功，Access-Accept"),
    ("Accounting",      "radacct 有 Start record，AcctStopTime NULL"),
    ("daloRADIUS Dashboard","線上使用者、accounting 報表正常顯示"),
    ("使用者管理",      "daloRADIUS 新增帳號 → 可正常登入 portal"),
    ("Portal 客製",     "客製 HTML + logo，client 看到新介面"),
    ("WISPr 限速",      "1Mbps → wget 1MB 實測 8.0s（~125 KB/s）"),
    ("Firewall Deny",   "未認證 client 送 SSH port 22 → nftables 擋掉"),
    ("CoA 踢人",        "Disconnect-ACK 收到，acctcause=Admin-Reset"),
]

col_x3 = [Inches(0.3), Inches(6.8)]
for i, (item, desc) in enumerate(tests):
    col = i // 6
    row = i % 6
    x = col_x3[col]
    y = Inches(1.6) + row * Inches(0.92)
    bg = RGBColor(0xEA,0xFA,0xEA) if i < 11 else C_LGRAY
    add_rect(s, x, y, Inches(6.2), Inches(0.82), fill=bg)
    add_rect(s, x, y, Inches(0.55), Inches(0.82), fill=C_GREEN)
    add_text(s, "✔", x+Inches(0.08), y+Inches(0.18), Inches(0.42), Inches(0.42),
             size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, item, x+Inches(0.65), y+Inches(0.06), Inches(5.4), Inches(0.38),
             size=12, bold=True, color=C_NAVY)
    add_text(s, desc, x+Inches(0.65), y+Inches(0.44), Inches(5.4), Inches(0.32),
             size=10, color=C_MGRAY)

add_text(s, "★  v1 MVP 全數通過 — 可示範、可交付",
         Inches(0.3), Inches(7.08), Inches(12.7), Inches(0.34),
         size=14, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — v2 規劃
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
header_bar(s, "後續 v2 規劃", "Cellular WAN + 壓測 + 監控 + 安全加固")
orange_accent(s)

v2_items = [
    (C_TEAL,   "📡  Cellular WAN",
     ["mmcli -L 列舉 modem", "mmcli -m 0 --simple-connect apn=internet",
      "systemd unit 監控 wwan0 斷線重連", "08-cellular.sh（目前為 stub）"]),
    (C_NAVY,   "🔀  WAN Failover",
     ["keepalived + track_script ping 偵測", "eth0 斷 → wwan0 接管 default route",
      "切換時 conntrack -F flush", "failover < 30s 目標"]),
    (C_ORANGE, "🔬  壓力測試",
     ["100 client 同時登入", "iperf3 throughput ≥ 1 Gbps（有線）",
      "72 hr 連續跑 RAM / CPU / conntrack", "radacct 資料完整性驗證"]),
    (C_GREEN,  "🔒  安全加固",
     ["daloRADIUS admin 密碼改掉（現在是 radius）", "self-signed → Let's Encrypt",
      "SNMP v3 createUser / rouser", "rsyslog 遠端轉送 @@<syslog-server>:514"]),
]

for i, (bg, title, items) in enumerate(v2_items):
    col = i % 2
    row = i // 2
    x = Inches(0.3) + col * Inches(6.5)
    y = Inches(1.6) + row * Inches(2.75)
    add_rect(s, x, y, Inches(6.2), Inches(2.6), fill=C_LGRAY)
    add_rect(s, x, y, Inches(6.2), Inches(0.5), fill=bg)
    add_text(s, title, x+Inches(0.15), y+Inches(0.06), Inches(5.9), Inches(0.42),
             size=14, bold=True, color=C_WHITE)
    tp4 = y + Inches(0.65)
    for item in items:
        add_text(s, "▸  " + item, x+Inches(0.2), tp4, Inches(5.8), Inches(0.4),
                 size=11.5, color=C_DKGRAY)
        tp4 += Inches(0.44)

# ════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════
import os
out_repo    = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CaptivePortal_FAE_Guide.pptx")
out_desktop = r"C:\Users\andy_chu\Desktop\CaptivePortal_FAE_Guide.pptx"
prs.save(out_repo)
print(f"Saved: {out_repo}")
try:
    prs.save(out_desktop)
    print(f"Saved: {out_desktop}")
except Exception as e:
    print(f"Desktop save skipped: {e}")
print(f"Slides: {len(prs.slides)}")
